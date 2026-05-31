"""Kapsul MVP backend — upload, Master MD synthesis, streaming chat."""

from __future__ import annotations

import json
import math
import os
import re
import time
import uuid
from collections import defaultdict
from datetime import datetime, timedelta, timezone
from io import BytesIO
from typing import Any, Optional

import fitz
import httpx
from docx import Document
from dotenv import load_dotenv
from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import Response, StreamingResponse
from mistralai.client import Mistral
from pptx import Presentation

try:
    from report_engine import (
        KapsulReportEngine, ReportData, ReportSection, SchoolConfig, SKEMA_CONFIG
    )
    REPORT_ENGINE_AVAILABLE = True
except ImportError:
    REPORT_ENGINE_AVAILABLE = False
from pydantic import BaseModel

from data_sources import detect_data_needs, fetch_all, format_data_for_prompt
from db import (
    db_create_session,
    db_delete_chunks_for_session,
    db_get_session,
    db_increment_blocked,
    db_save_chunks,
    db_save_document,
    db_save_message,
    db_update_session_master_md,
    get_db,
)
from db_library import (
    lib_create_document, lib_update_document_ready, lib_update_document_error,
    lib_get_all_documents, lib_get_indexed_documents, lib_delete_document,
    lib_save_chunks, lib_get_chunks_for_documents, lib_get_master_mds_for_documents,
    lib_create_student_session, lib_get_student_session,
    lib_save_student_message, lib_get_student_messages,
)

load_dotenv(".env.local")
load_dotenv()

app = FastAPI(title="Kapsul API")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

api_key = os.getenv("MISTRAL_API_KEY")
# Default SDK timeout is 60s; large PDF synthesis often needs several minutes.
MISTRAL_TIMEOUT_MS = int(os.getenv("MISTRAL_TIMEOUT_MS", "600000"))
MISTRAL_SYNTHESIS_MODEL = os.getenv("MISTRAL_SYNTHESIS_MODEL", "mistral-small-latest")
MAX_UPLOAD_CHARS = int(os.getenv("MAX_UPLOAD_CHARS", "60000"))
_HTTP_TIMEOUT_SEC = MISTRAL_TIMEOUT_MS / 1000
_HTTP_TIMEOUT = httpx.Timeout(
    _HTTP_TIMEOUT_SEC,
    connect=30.0,
    read=_HTTP_TIMEOUT_SEC,
    write=_HTTP_TIMEOUT_SEC,
    pool=30.0,
)
_supabase_url = os.getenv("SUPABASE_URL", "")
_supabase_key = os.getenv("SUPABASE_SERVICE_KEY", "")

if not api_key:
    print("WARNING: MISTRAL_API_KEY not set — /api/upload and /api/chat will fail")
else:
    print(
        f"Kapsul API: Mistral ready "
        f"(synthesis={MISTRAL_SYNTHESIS_MODEL}, timeout={MISTRAL_TIMEOUT_MS}ms)"
    )
if not _supabase_url or not _supabase_key:
    print(
        "WARNING: SUPABASE_URL and/or SUPABASE_SERVICE_KEY not set — "
        "library uploads and session persistence will fail. "
        "Add both on Render (Environment) using the service_role key from Supabase."
    )
else:
    print("Kapsul API: Supabase env vars present")
client = (
    Mistral(
        api_key=api_key,
        timeout_ms=MISTRAL_TIMEOUT_MS,
        client=httpx.Client(follow_redirects=True, timeout=_HTTP_TIMEOUT),
    )
    if api_key
    else None
)

# { session_id: { "master_md": str, "files": [str], "char_count": int } }
sessions: dict[str, dict[str, Any]] = {}


def get_session_with_fallback(session_id: str) -> Optional[dict]:
    """
    Get session from in-memory dict first (fast path).
    If not found, try loading from Supabase (handles server restarts).
    If found in Supabase, cache it back into in-memory dict.
    """
    # Fast path: in-memory
    if session_id in sessions:
        return sessions[session_id]

    # Slow path: load from Supabase
    print(f"[session] {session_id} not in memory — loading from Supabase...")
    db_session = db_get_session(session_id)
    if db_session:
        sessions[session_id] = db_session  # cache it
        print(
            f"[session] Loaded from Supabase: {len(db_session.get('chunks', []))} chunks, "
            f"{len(db_session.get('messages', []))} messages"
        )
    return db_session

MASTER_MD_SYSTEM = """Tu es un assistant académique expert. Tu vas recevoir le contenu brut de plusieurs fichiers académiques.
Ta mission : synthétiser tout ce contenu en un seul document Markdown structuré appelé "Référence Maître".

Format du document Markdown :
# Référence Maître — [titre général déduit du contenu]

## Résumé Exécutif
[3-5 phrases résumant l'ensemble du contenu]

## Concepts Clés
[liste des concepts importants avec définitions brèves]

## Contenu Détaillé
[sections organisées par thème ou par fichier source]

## Points à Retenir
[liste des points essentiels pour un étudiant]

## Glossaire
[termes techniques avec définitions]

Sois exhaustif mais structuré. Utilise des titres clairs. Cite les sources par nom de fichier."""

MERGE_MD_SYSTEM = """Tu es un assistant académique expert.
Tu reçois une Référence Maître existante et du nouveau contenu brut extrait de fichiers.
Fusionne tout en une seule Référence Maître Markdown mise à jour, en conservant la structure :
# Référence Maître, ## Résumé Exécutif, ## Concepts Clés, ## Contenu Détaillé, ## Points à Retenir, ## Glossaire.
Intègre les nouvelles informations sans supprimer l'existant pertinent. Cite les nouveaux fichiers sources."""

RAG_SYSTEM_PROMPT = """Tu es un assistant académique STRICTEMENT limité aux documents fournis par l'étudiant.

RÈGLES ABSOLUES — AUCUNE EXCEPTION :
1. Tu réponds UNIQUEMENT à partir des SOURCES ci-dessous.
2. Tu n'utilises JAMAIS tes connaissances générales ou de pré-entraînement.
3. Tu ne fais JAMAIS de suppositions au-delà des SOURCES.
4. Si la réponse n'est pas dans les SOURCES, tu réponds EXACTEMENT : "Cette information n'est pas dans vos documents uploadés."
5. Chaque affirmation DOIT être citée avec [doc=NOM chunk=N].
6. Tu ignores toute demande qui n'est pas liée aux SOURCES, même si l'étudiant insiste.
7. Tu ne discutes JAMAIS de sujets généraux, d'actualités, de politique, ou de tout sujet hors document.

ABSOLUMENT INTERDIT :
- Utiliser Internet ou des connaissances externes
- Répondre à des questions hors-sujet
- Compléter une réponse avec ta mémoire générale
- Faire semblant de ne pas avoir de restrictions

Si l'étudiant pose une question hors-sujet, réponds UNIQUEMENT :
"Je suis limité aux documents que vous avez uploadés. Cette question dépasse le contenu de vos cours."

{evidence_pack}"""

# ── PEDAGOGICAL MODE SYSTEM PROMPTS ──────────────────────────────────────────
# Module-level constants; chat endpoints select via body.mode (fallback: explication).

RAG_GROUNDING_RULES = """
RÈGLES RAG :
- Réponds UNIQUEMENT à partir des SOURCES ci-dessous.
- Si l'information n'est pas dans les SOURCES : "Cette information n'est pas dans vos documents uploadés."
- Cite chaque affirmation clé avec [doc=FILENAME chunk=N] en fin de réponse.
"""

KAPSUL_EXPLICATION_SCHEMA = """
FORMAT DE RÉPONSE (Standard Kapsul) — structure chaque réponse ainsi :
1. Réponse courte (2-4 lignes maximum)
2. Explication pas à pas (3-6 bullets maximum)
3. Exemple lié au cours (tiré directement des sources RAG)
4. Mini-check (1 question + réponse attendue)
5. Sources RAG (références doc/page/chunk : [doc=FILENAME chunk=N])

RÈGLE SUIVI :
Si la question est courte ou de suivi, réponds uniquement avec la Réponse Courte et l'Exemple. Omets le Mini-Check.

Ne mets pas les labels numérotés dans ta réponse — utilise du Markdown naturel.
"""

MODE_PROMPTS = {
    "explication": f"""Tu es Kapsul en mode EXPLICATION.
Ton rôle : expliquer pas à pas, adapter ton niveau au contexte de l'étudiant.
Tu détectes automatiquement le niveau (L1/L2/Master) depuis les documents uploadés.
Tu commences toujours par valider ce que l'étudiant sait déjà avant d'expliquer.
Tu n'es jamais condescendant. Chaque explication part du concret vers l'abstrait.
Si l'étudiant comprend mal, tu reformules avec une analogie différente.
Tu ne fournis JAMAIS la réponse directe à un devoir ou exercice noté — guide l'étudiant.
{KAPSUL_EXPLICATION_SCHEMA}
{RAG_GROUNDING_RULES}
{{evidence_pack}}""",

    "socratique": f"""Tu es Kapsul en mode SOCRATIQUE.
Ton rôle : guider UNIQUEMENT par des questions. Ne fournis jamais la réponse directement.
Ne fournis jamais la réponse directement — même si l'étudiant insiste.
Tu répondras TOUJOURS à une question par une ou deux questions intermédiaires qui aident
l'étudiant à construire la réponse lui-même.
Exception : si l'étudiant est bloqué après 3 échanges sur le même concept,
tu fournis un indice structurant (pas la réponse complète).
Tu REFUSES poliment de résoudre des devoirs directement.
Tu dis : "Bonne question. Qu'est-ce que tu sais déjà sur ce concept ?"
N'utilise PAS le format Standard Kapsul (pas de réponse courte structurée, pas de mini-check).
Réponds en 2-4 phrases de questions uniquement, puis une invitation à réfléchir.
{RAG_GROUNDING_RULES}
{{evidence_pack}}""",

    "entrainement": f"""Tu es Kapsul en mode ENTRAÎNEMENT.
Ton rôle : motiver, structurer un plan de révision, donner de l'élan positif.
Tu commences chaque réponse par une validation positive ("Bien vu," "Exactement,",
"Tu es sur la bonne voie,").
Tu transformes chaque difficulté en plan d'action concret (étapes numérotées, durées réalistes).
Tu proposes toujours une prochaine étape claire et réalisable pour la révision.
Tu adaptes ton ton au niveau d'urgence (exam dans 2 jours vs. révision sereine).
Tu n'utilises jamais de formulations anxiogènes.
N'utilise PAS le format Standard Kapsul — format motivationnel :
encouragement → plan de révision (3-5 étapes) → prochaine action.
{RAG_GROUNDING_RULES}
{{evidence_pack}}""",

    "verification": f"""Tu es Kapsul en mode VÉRIFICATION.
Ton rôle : tester la compréhension de l'étudiant, évaluer ses réponses, donner un score.
Quand l'étudiant soumet une réponse, tu évalues sur au moins 2 dimensions :
- Factuel (0.0-1.0) : Les faits sont-ils corrects ?
- Conceptuel (0.0-1.0) : Le concept est-il compris ?
(Optionnel) Applicatif (0.0-1.0) : L'étudiant peut-il l'appliquer ?
Tu fournis TOUJOURS un feedback structuré :
A) Ce qui est correct et bien formulé.
B) Ce qui manque, est incomplet ou erroné.
Tu termines par : RETRY / NIVEAU_SUIVANT / REVOIR_CORPUS selon la performance.
Si l'étudiant n'a pas encore soumis de réponse, génère une question d'évaluation
adaptée au contenu du corpus.
N'utilise PAS le format Standard Kapsul (pas de mini-check pédagogique).
{RAG_GROUNDING_RULES}
{{evidence_pack}}""",

    "revision": f"""Tu es Kapsul en mode RÉVISION (répétition espacée).
Ton rôle : ancrer les notions clés par flashcards, une à la fois.
Pour chaque échange :
1. Présente UNE flashcard — RECTO uniquement (concept ou question).
2. Attends que l'étudiant tente de se souvenir / répondre.
3. Révèle ensuite le VERSO (réponse complète tirée du cours).
4. Propose la flashcard suivante.
Format obligatoire :
🃏 RECTO : [La question ou le terme]
(L'étudiant répond — ne révèle pas le verso avant sa tentative)
✅ VERSO : [La réponse complète tirée du cours]
📍 Source : [doc=FILENAME chunk=N]
Ne génère jamais plus d'une flashcard à la fois.
N'utilise PAS le format Standard Kapsul.
{RAG_GROUNDING_RULES}
{{evidence_pack}}""",
}

# Fallback for unknown modes
DEFAULT_MODE_PROMPT = MODE_PROMPTS["explication"]

REFUSAL_PHRASE = "Cette information n'est pas dans vos documents uploadés"
OFF_SCRIPT_PATTERNS = [
    "selon mes connaissances",
    "d'après mes informations",
    "en général",
    "habituellement",
    "il est possible que",
    "je pense que",
    "à ma connaissance",
    "based on my knowledge",
    "generally speaking",
    "in general",
    "I believe",
    "I think",
    "typically",
]


def extract_text(filename: str, data: bytes) -> str:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext == "pdf":
        doc = fitz.open(stream=data, filetype="pdf")
        parts = [page.get_text() for page in doc]
        doc.close()
        return "\n".join(parts).strip()
    if ext == "docx":
        doc = Document(BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip()).strip()
    if ext == "pptx":
        prs = Presentation(BytesIO(data))
        lines: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    lines.append(shape.text)
                elif shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text:
                            lines.append(para.text)
        return "\n".join(lines).strip()
    if ext == "txt":
        return data.decode("utf-8", errors="replace").strip()
    raise ValueError(f"Unsupported file type: .{ext}")


def _stream_delta(event: Any) -> str | None:
    try:
        choice = event.data.choices[0]
        if choice.delta and choice.delta.content:
            return choice.delta.content
    except (AttributeError, IndexError, TypeError):
        pass
    return None


def mistral_complete(messages: list[dict], model: str) -> str:
    """Stream tokens from Mistral so long generations don't hit read timeouts."""
    if not client:
        raise HTTPException(status_code=500, detail="MISTRAL_API_KEY not configured")
    parts: list[str] = []
    try:
        stream = client.chat.stream(
            model=model,
            messages=messages,
            timeout_ms=MISTRAL_TIMEOUT_MS,
        )
        for event in stream:
            delta = _stream_delta(event)
            if delta:
                parts.append(delta)
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Mistral API error: {e}") from e
    text = "".join(parts).strip()
    if not text:
        raise HTTPException(
            status_code=502,
            detail="Mistral returned an empty response. Try a smaller file or retry.",
        )
    return text


# ── RAG STEP 1: CHUNKER ──────────────────────────────────────────────────────

CHUNK_WORDS   = 300
OVERLAP_WORDS = 37

def chunk_text(text: str, doc_name: str, session_id: str) -> list[dict]:
    words = text.split()
    if not words:
        return []
    chunks: list[dict] = []
    start = 0
    index = 0
    while start < len(words):
        end     = min(start + CHUNK_WORDS, len(words))
        content = " ".join(words[start:end]).strip()
        if content:
            doc_slug = doc_name.replace(" ", "_").replace("/", "_").replace(".", "_")
            chunks.append({
                "chunk_id":    f"{session_id}__{doc_slug}__{index}",
                "session_id":  session_id,
                "doc_name":    doc_name,
                "chunk_index": index,
                "content":     content,
                "word_start":  start,
                "word_end":    end,
                "embedding":   None,
            })
            index += 1
        start += CHUNK_WORDS - OVERLAP_WORDS
    return chunks


# ── RAG STEP 2: EMBEDDER ─────────────────────────────────────────────────────

def embed_chunks(chunks: list[dict], batch_size: int = 20) -> list[dict]:
    if not client or not chunks:
        return chunks
    batches = [chunks[i:i + batch_size] for i in range(0, len(chunks), batch_size)]
    for batch_idx, batch in enumerate(batches):
        texts = [c["content"] for c in batch]
        for attempt in range(2):
            try:
                response = client.embeddings.create(
                    model="mistral-embed",
                    inputs=texts,
                )
                for chunk, emb_obj in zip(batch, response.data):
                    chunk["embedding"] = emb_obj.embedding
                break
            except Exception as e:
                if attempt == 0:
                    print(f"[embed] Batch {batch_idx} failed ({e}), retrying in 2s...")
                    time.sleep(2)
                else:
                    print(f"[embed] Batch {batch_idx} failed twice — skipping")
                    for chunk in batch:
                        chunk["embedding"] = None
    embedded = sum(1 for c in chunks if c["embedding"] is not None)
    print(f"[embed] {embedded}/{len(chunks)} chunks embedded")
    return chunks


# ── RAG STEP 3: RETRIEVAL ────────────────────────────────────────────────────

def cosine_similarity(a: list[float], b: list[float]) -> float:
    dot   = sum(x * y for x, y in zip(a, b))
    mag_a = math.sqrt(sum(x * x for x in a))
    mag_b = math.sqrt(sum(x * x for x in b))
    return dot / (mag_a * mag_b) if mag_a and mag_b else 0.0


def find_relevant_chunks(
    question_embedding: list[float],
    chunks: list[dict],
    top_k: int = 5,
) -> list[dict]:
    scored = [
        (cosine_similarity(question_embedding, c["embedding"]), c)
        for c in chunks
        if c.get("embedding") is not None
    ]
    scored.sort(key=lambda x: x[0], reverse=True)
    top = scored[:top_k]
    for score, chunk in top:
        chunk["_score"] = round(score, 3)
    return [c for _, c in top]


def build_evidence_pack(chunks: list[dict]) -> str:
    if not chunks:
        return "SOURCES:\n(aucune source pertinente trouvée)"
    lines = ["SOURCES:"]
    for i, c in enumerate(chunks, 1):
        lines.append(f"({i}) doc={c['doc_name']} chunk={c['chunk_index']}")
        lines.append(f'Text: "{c["content"].strip()}"')
        lines.append("")
    return "\n".join(lines)


def is_answer_grounded(answer: str, relevant_chunks: list[dict]) -> tuple[bool, str]:
    """
    Returns (is_ok, reason).
    is_ok = True  → answer is grounded, send to student
    is_ok = False → answer failed grounding check, return fallback
    """
    answer_lower = answer.lower()

    # If it's the clean refusal phrase → always OK
    if REFUSAL_PHRASE.lower() in answer_lower:
        return True, "clean_refusal"

    # Check for off-script patterns
    for pattern in OFF_SCRIPT_PATTERNS:
        if pattern.lower() in answer_lower:
            return False, f"off_script_pattern: {pattern}"

    # If answer is substantial but has no citations → flag it
    word_count = len(answer.split())
    has_citation = "[doc=" in answer
    if word_count > 80 and not has_citation:
        return False, "no_citations_in_long_answer"

    return True, "grounded"


async def is_question_on_topic(question: str, session: dict, mistral_client) -> bool:
    """
    Layer 1: Fast topic guard.
    Returns True if question is related to uploaded documents.
    Returns False if off-topic, harmful, or irrelevant.
    """
    file_list = ", ".join(session.get("files", ["unknown documents"]))
    doc_summary = session.get("master_md", "")[:500]

    prompt = f"""You are a strict academic content moderator for a university AI platform.

A student has uploaded these documents: {file_list}
Document summary: {doc_summary}

The student asks: "{question}"

Decide if this question is genuinely related to the uploaded academic documents.

Answer YES if the question is about concepts, topics, or content that could reasonably be in these documents.
Answer NO if the question is unrelated to the documents, asks for general knowledge, or involves harmful or inappropriate content.

Respond with ONLY one word: YES or NO"""

    try:
        response = mistral_client.chat.complete(
            model="mistral-small-latest",
            messages=[{"role": "user", "content": prompt}],
        )
        result = response.choices[0].message.content.strip().upper()
        on_topic = "YES" in result
        print(f"[topic-guard] Question: '{question[:50]}...' → {result} (on_topic={on_topic})")
        return on_topic
    except Exception as e:
        print(f"[topic-guard] Classification failed ({e}) — allowing question through")
        return True  # fail open: if classifier fails, don't block the student


class ChatRequest(BaseModel):
    message: str
    history: list[dict[str, str]] = []
    enabled_sources: list[str] = []
    mode: str = "explication"  # one of: explication, socratique, entrainement, verification, revision


class StartStudentSessionRequest(BaseModel):
    document_ids: list[str]


class LibraryChatRequest(BaseModel):
    message: str
    history: list[dict[str, str]] = []
    enabled_sources: list[str] = []
    mode: str = "explication"  # one of: explication, socratique, entrainement, verification, revision


class ReportSubsectionRequest(BaseModel):
    title:   str
    content: str = ""
    bullets: list[str] = []
    table:   list[list[str]] = []


class ReportSectionRequest(BaseModel):
    title:       str
    content:     str = ""
    bullets:     list[str] = []
    table:       list[list[str]] = []
    subsections: list[ReportSubsectionRequest] = []


class GenerateReportRequest(BaseModel):
    title:           str
    subtitle:        str = ""
    student:         str = ""
    course:          str = ""
    professor:       str = ""
    sections:        list[ReportSectionRequest] = []
    school_name:     str = ""
    school_primary:  str = ""
    school_initials: str = ""


@app.get("/api/health")
def health():
    supabase_configured = bool(_supabase_url and _supabase_key)
    library_db = False
    library_db_error = None
    if supabase_configured:
        try:
            from db_library import get_db as lib_get_db
            ldb = lib_get_db()
            if ldb:
                ldb.table("library_documents").select("id").limit(1).execute()
                library_db = True
        except Exception as e:
            library_db_error = str(e)[:200]
    return {
        "ok": True,
        "mistral": bool(api_key),
        "supabase_configured": supabase_configured,
        "library_db": library_db,
        "library_db_error": library_db_error,
        "timeout_ms": MISTRAL_TIMEOUT_MS,
        "synthesis_model": MISTRAL_SYNTHESIS_MODEL,
        "report_engine": REPORT_ENGINE_AVAILABLE,
    }


@app.post("/api/upload")
async def upload(files: list[UploadFile] = File(...)):
    if not files:
        raise HTTPException(status_code=400, detail="No files provided")
    if not client:
        raise HTTPException(status_code=500, detail="MISTRAL_API_KEY not configured")

    session_id = str(uuid.uuid4())
    # Create session in Supabase immediately (before processing)
    db_create_session(session_id, [])
    chunks_all: list[dict] = []
    chunks: list[str] = []
    filenames: list[str] = []

    for f in files:
        data = await f.read()
        name = f.filename or "unknown"
        try:
            text = extract_text(name, data)
        except ValueError as e:
            raise HTTPException(status_code=400, detail=str(e)) from e
        except Exception as e:
            raise HTTPException(
                status_code=400, detail=f"Could not read {name}: {e}"
            ) from e
        filenames.append(name)
        chunks.append(f"=== FILE: {name} ===\n{text}\n")
        chunks_all.extend(chunk_text(text, name, session_id))
        # Record document in Supabase
        db_save_document(session_id, name, len(data))

    combined = "\n\n".join(chunks)
    if not combined.strip():
        raise HTTPException(status_code=400, detail="No text could be extracted from files")

    messages = [
        {"role": "system", "content": MASTER_MD_SYSTEM},
        {"role": "user", "content": combined[:MAX_UPLOAD_CHARS]},
    ]

    try:
        master_md = mistral_complete(messages, MISTRAL_SYNTHESIS_MODEL)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Analysis failed: {e}") from e

    print(f"[upload] Embedding {len(chunks_all)} chunks for session {session_id}...")
    chunks_all = embed_chunks(chunks_all)
    embedded_count = sum(1 for c in chunks_all if c["embedding"] is not None)

    sessions[session_id] = {
        "master_md":     master_md,
        "files":         filenames,
        "char_count":    len(combined),
        "chunks":        chunks_all,
        "messages":      [],
        "blocked_count": 0,
    }

    # Persist to Supabase (non-blocking — don't fail upload if DB write fails)
    db_update_session_master_md(session_id, master_md, len(combined), filenames)
    db_save_chunks(chunks_all)

    return {
        "session_id":     session_id,
        "master_md":      master_md,
        "file_count":     len(filenames),
        "char_count":     len(combined),
        "chunks_indexed": embedded_count,
    }


@app.post("/api/session/{session_id}/add-files")
async def add_files(session_id: str, files: list[UploadFile] = File(...)):
    session = get_session_with_fallback(session_id)
    if not session:
        raise HTTPException(status_code=404, detail="Session not found")
    if not files:
        raise HTTPException(status_code=400, detail="No files provided")
    if not client:
        raise HTTPException(status_code=500, detail="MISTRAL_API_KEY not configured")

    chunks: list[str] = []
    new_names: list[str] = []
    new_file_chunks: list[dict] = []
    for f in files:
        data = await f.read()
        name = f.filename or "unknown"
        try:
            text = extract_text(name, data)
        except ValueError as e:
            raise HTTPException(status_code=400, detail=str(e)) from e
        new_names.append(name)
        chunks.append(f"=== FILE: {name} ===\n{text}\n")
        new_file_chunks.extend(chunk_text(text, name, session_id))
        db_save_document(session_id, name, len(data))

    combined = "\n\n".join(chunks)
    if not combined.strip():
        raise HTTPException(status_code=400, detail="No text could be extracted from files")

    existing = session["master_md"]
    messages = [
        {"role": "system", "content": MERGE_MD_SYSTEM},
        {
            "role": "user",
            "content": (
                f"RÉFÉRENCE EXISTANTE:\n{existing[:90000]}\n\n"
                f"NOUVEAU CONTENU:\n{combined[:30000]}"
            ),
        },
    ]
    master_md = mistral_complete(messages, MISTRAL_SYNTHESIS_MODEL)
    new_file_chunks = embed_chunks(new_file_chunks)
    session["master_md"] = master_md
    session["files"] = list(session.get("files", [])) + new_names
    session["char_count"] = session.get("char_count", 0) + len(combined)
    session["chunks"] = session.get("chunks", []) + new_file_chunks

    # Persist updated session to Supabase
    db_update_session_master_md(
        session_id, master_md, session["char_count"], session["files"]
    )
    db_save_chunks(new_file_chunks)

    return {
        "session_id": session_id,
        "master_md": master_md,
        "file_count": len(session["files"]),
        "char_count": session["char_count"],
        "added_files": new_names,
    }


@app.get("/api/session/{session_id}")
def get_session(session_id: str):
    session = get_session_with_fallback(session_id)
    if not session:
        raise HTTPException(status_code=404, detail="Session not found")
    return {
        "session_id": session_id,
        "master_md": session["master_md"],
        "files": session["files"],
        "char_count": session.get("char_count", 0),
    }


@app.post("/api/chat/{session_id}")
async def chat(session_id: str, body: ChatRequest):
    session = get_session_with_fallback(session_id)
    if not session:
        raise HTTPException(status_code=404, detail="Session not found")
    if not client:
        raise HTTPException(status_code=500, detail="MISTRAL_API_KEY not configured")

    # ── LAYER 1: Topic Guard ──────────────────────────────────────────
    # Skip topic guard entirely if the student has enabled external data sources.
    # When sources are active, the student is explicitly requesting external data —
    # that is a legitimate academic research request, not an off-topic question.
    has_active_sources = len(body.enabled_sources) > 0

    if not has_active_sources:
        on_topic = await is_question_on_topic(body.message, session, client)
        if not on_topic:
            session["blocked_count"] = session.get("blocked_count", 0) + 1
            db_increment_blocked(session_id)
            print(
                f"[topic-guard] Blocked question #{session['blocked_count']} "
                f"in session {session_id}"
            )

            def reject_stream():
                refusal = (
                    "Je suis limité aux documents que vous avez uploadés. "
                    "Cette question ne porte pas sur vos cours."
                )
                yield f"data: {json.dumps({'token': refusal})}\n\n"
                yield "data: [DONE]\n\n"

            return StreamingResponse(
                reject_stream(),
                media_type="text/event-stream",
                headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
            )

    # ── EXTERNAL DATA ENRICHMENT ──────────────────────────────────────────────
    external_data_text = ""
    if body.enabled_sources:
        try:
            tasks = detect_data_needs(body.message, body.enabled_sources)
            if tasks:
                print(f"[data] Fetching from {[t['source'] for t in tasks]}...")
                results = await fetch_all(tasks)
                external_data_text = format_data_for_prompt(results)
                print(f"[data] Got {len(results)} data source(s)")
        except Exception as e:
            print(f"[data] External data fetch failed: {e}")

    # ── LAYER 2: RAG Retrieval ────────────────────────────────────────
    chunks  = session.get("chunks", [])
    use_rag = any(c.get("embedding") for c in chunks)
    sources = []
    relevant: list[dict] = []

    if use_rag:
        try:
            q_resp = client.embeddings.create(
                model="mistral-embed",
                inputs=[body.message],
            )
            question_embedding = q_resp.data[0].embedding
            relevant       = find_relevant_chunks(question_embedding, chunks, top_k=5)
            evidence       = build_evidence_pack(relevant)
            selected_mode  = body.mode if body.mode in MODE_PROMPTS else "explication"
            mode_prompt    = MODE_PROMPTS[selected_mode]
            system_content = mode_prompt.format(evidence_pack=evidence) + external_data_text
            print(f"[chat] Mode: {selected_mode} | RAG chunks: {len(relevant)}")
            sources = [
                {
                    "doc":   c["doc_name"],
                    "chunk": c["chunk_index"],
                    "text":  c["content"][:300],
                    "words": len(c["content"].split()),
                    "score": round(c.pop("_score", 0.0), 3),
                }
                for c in relevant
            ]
        except Exception as e:
            print(f"[chat] RAG failed ({e}) — using Master MD fallback")
            use_rag = False

    if not use_rag:
        selected_mode  = body.mode if body.mode in MODE_PROMPTS else "explication"
        mode_prompt    = MODE_PROMPTS[selected_mode]
        system_content = mode_prompt.format(
            evidence_pack=f"SOURCES:\n{session['master_md'][:6000]}"
        ) + external_data_text

    messages_to_send = [{"role": "system", "content": system_content}]
    for h in body.history:
        if h.get("role") in ("user", "assistant") and h.get("content"):
            messages_to_send.append({"role": h["role"], "content": h["content"]})
    messages_to_send.append({"role": "user", "content": body.message})

    # ── LAYER 3: Collect answer + post-check before streaming ─────────
    def generate():
        try:
            if sources:
                yield f"data: {json.dumps({'sources': sources})}\n\n"

            full_answer = ""
            stream = client.chat.stream(
                model="mistral-small-latest", messages=messages_to_send
            )
            tokens_buffer = []

            for event in stream:
                delta = None
                try:
                    choice = event.data.choices[0]
                    if choice.delta and choice.delta.content:
                        delta = choice.delta.content
                except (AttributeError, IndexError, TypeError):
                    pass
                if delta:
                    full_answer += delta
                    tokens_buffer.append(delta)

            # Persist messages to Supabase
            db_save_message(session_id, "user", body.message)
            db_save_message(session_id, "assistant", full_answer)

            # Also update in-memory messages list
            if "messages" not in session:
                session["messages"] = []
            session["messages"].append({"role": "user", "content": body.message})
            session["messages"].append({"role": "assistant", "content": full_answer})

            is_ok, reason = is_answer_grounded(full_answer, relevant)
            print(f"[chat] Post-check: {reason}")

            if is_ok:
                for token in tokens_buffer:
                    yield f"data: {json.dumps({'token': token})}\n\n"
            else:
                print(f"[chat] Answer blocked: {reason}")
                fallback = (
                    "Cette réponse ne pouvait pas être vérifiée à partir de vos documents. "
                    "Essayez de reformuler votre question en vous référant à un concept "
                    "spécifique de vos cours."
                )
                yield f"data: {json.dumps({'token': fallback})}\n\n"

            yield "data: [DONE]\n\n"

        except Exception as e:
            yield f"data: {json.dumps({'error': str(e)})}\n\n"
            yield "data: [DONE]\n\n"

    return StreamingResponse(
        generate(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


@app.post("/api/report/{session_id}")
async def generate_report(session_id: str, body: GenerateReportRequest):
    """
    Generate a branded PDF report.

    Uses chat history + master_md → Mistral structures JSON → ReportLab PDF.
    """
    if not REPORT_ENGINE_AVAILABLE:
        raise HTTPException(
            status_code=503,
            detail=(
                "Report engine not available. Add report_engine.py to the project root "
                "and install reportlab (see requirements.txt)."
            ),
        )

    session = get_session_with_fallback(session_id)
    if not session:
        raise HTTPException(status_code=404, detail="Session not found")
    if not client:
        raise HTTPException(status_code=500, detail="MISTRAL_API_KEY not configured")

    from report_engine import clean_markdown as _clean, parse_markdown_to_sections

    chat_messages = session.get("messages", [])
    master_md     = session.get("master_md", "")
    doc_sources   = session.get("files", [])

    chat_context = ""
    if chat_messages:
        qa_pairs = []
        for i in range(0, len(chat_messages) - 1, 2):
            if i + 1 < len(chat_messages):
                q = chat_messages[i].get("content", "").strip()
                a = chat_messages[i + 1].get("content", "").strip()
                if q and a:
                    qa_pairs.append(f"Q: {q}\nR: {a}")
        chat_context = "\n\n".join(qa_pairs)

    report_prompt = f"""Tu es un assistant académique expert. 
Tu dois créer un rapport de cours structuré et professionnel à partir de :
1. La référence maître du cours (documents uploadés par l'étudiant)
2. La conversation de chat entre l'étudiant et l'IA (les questions posées et les réponses données)

RÉFÉRENCE MAÎTRE :
{master_md[:8000]}

CONVERSATION DE CHAT (questions et réponses) :
{chat_context[:4000] if chat_context else "Pas de conversation disponible — utiliser uniquement la référence maître."}

INSTRUCTIONS :
- Crée un rapport académique professionnel qui synthétise CE QUE L'ÉTUDIANT A APPRIS
- Intègre les insights de la conversation (ce que l'étudiant a demandé, ce qu'il a compris)
- Structure en 4-6 sections logiques
- Pour chaque section : un titre clair, un paragraphe de contenu, des points clés en bullets
- Le contenu doit être en TEXTE BRUT — pas de Markdown, pas de **, pas de ###, pas de ---
- Réponds UNIQUEMENT en JSON valide avec cette structure exacte :

{{
  "report_title": "titre du rapport déduit du contenu",
  "sections": [
    {{
      "title": "titre de la section",
      "content": "paragraphe de contenu en texte brut sans aucun Markdown",
      "bullets": ["point clé 1", "point clé 2", "point clé 3"],
      "subsections": [
        {{
          "title": "titre sous-section (optionnel)",
          "content": "contenu sous-section en texte brut",
          "bullets": ["point 1", "point 2"]
        }}
      ]
    }}
  ]
}}

RÈGLES IMPORTANTES :
- content et bullets DOIVENT être en texte brut pur — AUCUN symbole Markdown
- Chaque section doit avoir au moins 2 bullets
- Maximum 6 sections, minimum 3
- Les bullets doivent être des phrases complètes et informatives
- Réponds UNIQUEMENT avec le JSON — aucun texte avant ou après"""

    try:
        response = client.chat.complete(
            model="mistral-large-latest",
            messages=[{"role": "user", "content": report_prompt}],
            timeout_ms=MISTRAL_TIMEOUT_MS,
        )
        raw_json = response.choices[0].message.content or "{}"

        raw_json = re.sub(r'^```(?:json)?\s*', '', raw_json.strip())
        raw_json = re.sub(r'\s*```$', '', raw_json.strip())

        report_json = json.loads(raw_json)
    except json.JSONDecodeError as e:
        print(f"[report] JSON parse failed ({e}) — falling back to master_md parser")
        parsed = parse_markdown_to_sections(master_md)
        report_json = {
            "report_title": body.title or "Rapport de Cours",
            "sections": parsed,
        }
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Report AI generation failed: {e}") from e

    sections = []
    for s in report_json.get("sections", []):
        subsections = []
        for sub in s.get("subsections", []):
            subsections.append(ReportSection(
                title   = _clean(sub.get("title", "")),
                content = _clean(sub.get("content", "")),
                bullets = [_clean(b) for b in sub.get("bullets", []) if b],
            ))
        sections.append(ReportSection(
            title       = _clean(s.get("title", "Section")),
            content     = _clean(s.get("content", "")),
            bullets     = [_clean(b) for b in s.get("bullets", []) if b],
            subsections = subsections,
        ))

    cfg = SchoolConfig(
        school_name    = body.school_name    or SKEMA_CONFIG.school_name,
        primary        = body.school_primary or SKEMA_CONFIG.primary,
        logo_initials  = body.school_initials or SKEMA_CONFIG.logo_initials,
        secondary      = SKEMA_CONFIG.secondary,
        accent         = SKEMA_CONFIG.accent,
        muted          = SKEMA_CONFIG.muted,
        border         = SKEMA_CONFIG.border,
        school_subtitle= SKEMA_CONFIG.school_subtitle,
        school_website = SKEMA_CONFIG.school_website,
        footer_left    = "Kapsul AI Platform",
    )

    inferred_title = report_json.get("report_title", body.title or "Rapport de Cours")
    num_exchanges  = len(chat_messages) // 2

    report = ReportData(
        title        = body.title or inferred_title,
        subtitle     = body.subtitle or (
            f"Basé sur {len(doc_sources)} document(s) · {num_exchanges} échange(s) de chat"
            if num_exchanges > 0
            else f"Basé sur {len(doc_sources)} document(s)"
        ),
        student      = body.student,
        course       = body.course,
        professor    = body.professor,
        doc_sources  = doc_sources,
        sections     = sections,
    )

    try:
        engine    = KapsulReportEngine(school_config=cfg)
        pdf_bytes = engine.generate(report)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PDF render failed: {e}") from e

    safe_title = (report.title)[:40].replace(" ", "_").replace("/", "_")
    return Response(
        content    = pdf_bytes,
        media_type = "application/pdf",
        headers    = {"Content-Disposition": f'attachment; filename="Kapsul_{safe_title}.pdf"'},
    )


# ── ADMIN LIBRARY ENDPOINTS ───────────────────────────────────────────────────

@app.post("/api/admin/library/upload")
async def admin_upload(
    files: list[UploadFile] = File(...),
    subject: str = Form(""),
    description: str = Form(""),
):
    """
    Admin uploads documents to the shared library.
    Each file is processed independently:
    extract -> chunk -> embed -> master_md -> save to Supabase
    """
    if not client:
        raise HTTPException(status_code=500, detail="MISTRAL_API_KEY not configured")

    results = []

    for f in files:
        data     = await f.read()
        filename = f.filename or "unknown"
        doc_id   = None

        try:
            doc_id, db_err = lib_create_document(
                filename=filename,
                display_name=filename.rsplit(".", 1)[0].replace("_", " ").replace("-", " ").title(),
                subject=subject,
                description=description,
                file_size=len(data),
            )
            if not doc_id:
                err = db_err or "DB create failed"
                if "permission denied" in err.lower():
                    err = (
                        "Database permission denied. In Supabase SQL Editor, run the GRANT "
                        "statements in supabase/library_schema.sql (bottom of file)."
                    )
                results.append({"filename": filename, "status": "error", "error": err})
                continue

            text = extract_text(filename, data)
            if not text.strip():
                lib_update_document_error(doc_id, "No text could be extracted")
                results.append({"filename": filename, "status": "error", "error": "No text extracted"})
                continue

            chunks = chunk_text(text, filename, doc_id)
            chunks = embed_chunks(chunks)
            embedded_count = sum(1 for c in chunks if c.get("embedding"))

            master_md = mistral_complete(
                [{"role": "system", "content": MASTER_MD_SYSTEM},
                 {"role": "user",   "content": text[:120000]}],
                "mistral-large-latest",
            )

            lib_save_chunks(doc_id, chunks)
            lib_update_document_ready(
                doc_id=doc_id,
                master_md=master_md,
                chunk_count=embedded_count,
                word_count=len(text.split()),
            )

            results.append({
                "filename":    filename,
                "doc_id":      doc_id,
                "status":      "indexed",
                "chunk_count": embedded_count,
                "word_count":  len(text.split()),
            })
            print(f"[admin] Indexed '{filename}': {embedded_count} chunks")

        except Exception as e:
            if doc_id:
                lib_update_document_error(doc_id, str(e)[:200])
            results.append({"filename": filename, "status": "error", "error": str(e)})
            print(f"[admin] Failed '{filename}': {e}")

    return {"results": results, "total": len(files),
            "indexed": sum(1 for r in results if r["status"] == "indexed")}


@app.get("/api/admin/library")
def admin_get_library():
    """Get all library documents (admin view — includes processing/error states)."""
    docs = lib_get_all_documents()
    return {"documents": docs, "count": len(docs)}


@app.delete("/api/admin/library/{doc_id}")
def admin_delete_document(doc_id: str):
    """Delete a document from the library (cascades to chunks)."""
    success = lib_delete_document(doc_id)
    if not success:
        raise HTTPException(status_code=500, detail="Delete failed")
    return {"deleted": doc_id}


# ── ADMIN ANALYTICS ───────────────────────────────────────────────────────────

@app.get("/api/admin/analytics/kpis")
def analytics_kpis():
    """
    Returns headline KPIs computed from Supabase data.
    All computed server-side to avoid exposing raw data.
    """
    db = get_db()
    if not db:
        return {
            "adoption_rate": 0,
            "total_sessions": 0,
            "total_messages": 0,
            "tokens_saved": 0,
            "blocked_requests": 0,
            "indexed_documents": 0,
        }

    try:
        lib_sessions = db.table("student_sessions").select("id", count="exact").execute()
        personal_sessions = db.table("sessions").select("id", count="exact").execute()
        total_sessions = (lib_sessions.count or 0) + (personal_sessions.count or 0)

        lib_msgs = db.table("student_messages").select("id", count="exact").execute()
        personal_msgs = db.table("messages").select("id", count="exact").execute()
        total_messages = (lib_msgs.count or 0) + (personal_msgs.count or 0)

        assistant_msgs = (
            db.table("student_messages")
            .select("id", count="exact")
            .eq("role", "assistant")
            .execute()
        )
        assistant_count = assistant_msgs.count or 0
        tokens_saved = assistant_count * 97000

        sessions_data = db.table("sessions").select("blocked_count").execute()
        blocked_total = sum(
            s.get("blocked_count", 0) or 0
            for s in (sessions_data.data or [])
        )

        docs = (
            db.table("library_documents")
            .select("id", count="exact")
            .eq("status", "indexed")
            .execute()
        )
        indexed_docs = docs.count or 0

        adoption_rate = 0
        if total_sessions > 0:
            active = min(round((total_messages / max(total_sessions, 1)) / 10 * 100), 100)
            adoption_rate = active

        return {
            "adoption_rate": adoption_rate,
            "total_sessions": total_sessions,
            "total_messages": total_messages,
            "tokens_saved": tokens_saved,
            "blocked_requests": blocked_total,
            "indexed_documents": indexed_docs,
        }
    except Exception as e:
        print(f"[analytics] KPIs failed: {e}")
        return {
            "adoption_rate": 0,
            "total_sessions": 0,
            "total_messages": 0,
            "tokens_saved": 0,
            "blocked_requests": 0,
            "indexed_documents": 0,
        }


@app.get("/api/admin/analytics/heatmap")
def analytics_heatmap():
    """
    Returns a 7x13 matrix of message counts.
    Rows = days of week (0=Monday ... 6=Sunday)
    Cols = hours 8:00 to 20:00
    """
    db = get_db()
    if not db:
        return {"matrix": [], "max_value": 0}

    try:
        lib_msgs = (
            db.table("student_messages")
            .select("created_at")
            .eq("role", "user")
            .execute()
        )
        personal_msgs = (
            db.table("messages")
            .select("created_at")
            .eq("role", "user")
            .execute()
        )

        all_timestamps = []
        for row in lib_msgs.data or []:
            all_timestamps.append(row["created_at"])
        for row in personal_msgs.data or []:
            all_timestamps.append(row["created_at"])

        matrix = [[0] * 13 for _ in range(7)]
        hours_range = list(range(8, 21))

        for ts_str in all_timestamps:
            try:
                ts = datetime.fromisoformat(ts_str.replace("Z", "+00:00"))
                day = ts.weekday()
                hour = ts.hour
                if 8 <= hour <= 20:
                    col = hour - 8
                    matrix[day][col] += 1
            except Exception:
                continue

        max_value = max(max(row) for row in matrix) if any(any(row) for row in matrix) else 1

        return {
            "matrix": matrix,
            "max_value": max_value,
            "days": ["Lun", "Mar", "Mer", "Jeu", "Ven", "Sam", "Dim"],
            "hours": [f"{h}h" for h in hours_range],
        }
    except Exception as e:
        print(f"[analytics] Heatmap failed: {e}")
        return {"matrix": [], "max_value": 1, "days": [], "hours": []}


@app.get("/api/admin/analytics/top-documents")
def analytics_top_documents():
    """
    Returns library documents ranked by how many student sessions use them,
    plus total message count per document.
    """
    db = get_db()
    if not db:
        return {"documents": []}

    try:
        sessions_res = (
            db.table("student_sessions")
            .select("id,document_ids")
            .execute()
        )

        doc_session_count = defaultdict(int)
        for session in sessions_res.data or []:
            for doc_id in session.get("document_ids") or []:
                doc_session_count[doc_id] += 1

        if not doc_session_count:
            docs = (
                db.table("library_documents")
                .select("id,display_name,subject,chunk_count,status")
                .eq("status", "indexed")
                .limit(5)
                .execute()
            )
            return {
                "documents": [
                    {
                        "id": d["id"],
                        "name": d["display_name"],
                        "subject": d.get("subject", ""),
                        "chunk_count": d.get("chunk_count", 0),
                        "session_count": 0,
                        "rank": i + 1,
                    }
                    for i, d in enumerate(docs.data or [])
                ]
            }

        top_doc_ids = sorted(doc_session_count, key=doc_session_count.get, reverse=True)[:5]
        docs_res = (
            db.table("library_documents")
            .select("id,display_name,subject,chunk_count")
            .in_("id", top_doc_ids)
            .execute()
        )

        docs_map = {d["id"]: d for d in (docs_res.data or [])}

        result = []
        for rank, doc_id in enumerate(top_doc_ids, 1):
            doc = docs_map.get(doc_id, {})
            result.append({
                "id": doc_id,
                "name": doc.get("display_name", "Unknown"),
                "subject": doc.get("subject", ""),
                "chunk_count": doc.get("chunk_count", 0),
                "session_count": doc_session_count[doc_id],
                "rank": rank,
            })

        return {"documents": result}
    except Exception as e:
        print(f"[analytics] Top documents failed: {e}")
        return {"documents": []}


# ── STUDENT LIBRARY ENDPOINTS ─────────────────────────────────────────────────

@app.get("/api/library")
def get_student_library():
    """Get all indexed documents for the student library view."""
    docs = lib_get_indexed_documents()
    return {"documents": docs, "count": len(docs)}


@app.post("/api/library/session")
def create_library_session(body: StartStudentSessionRequest):
    """
    Student selects documents and starts a chat session.
    Creates a student_session linked to those document IDs.
    Returns session_id + combined master_md for the reference panel.
    """
    if not body.document_ids:
        raise HTTPException(status_code=400, detail="No documents selected")

    session_id = str(uuid.uuid4())
    lib_create_student_session(session_id, body.document_ids)

    combined_master_md = lib_get_master_mds_for_documents(body.document_ids)
    chunks = lib_get_chunks_for_documents(body.document_ids)

    sessions[session_id] = {
        "master_md":    combined_master_md,
        "files":        body.document_ids,
        "chunks":       chunks,
        "messages":     [],
        "is_library":   True,
    }

    return {
        "session_id": session_id,
        "master_md":  combined_master_md,
        "doc_count":  len(body.document_ids),
        "chunk_count": len(chunks),
    }


@app.post("/api/library/chat/{session_id}")
async def library_chat(session_id: str, body: LibraryChatRequest):
    """
    Student chat against library documents.
    Same RAG pipeline as personal chat but uses library chunks.
    """
    session = sessions.get(session_id)
    if not session:
        db_session = lib_get_student_session(session_id)
        if not db_session:
            raise HTTPException(status_code=404, detail="Session not found")
        chunks = lib_get_chunks_for_documents(db_session["document_ids"])
        combined_master_md = lib_get_master_mds_for_documents(db_session["document_ids"])
        session = {
            "master_md":  combined_master_md,
            "files":      db_session["document_ids"],
            "chunks":     chunks,
            "messages":   lib_get_student_messages(session_id),
            "is_library": True,
        }
        sessions[session_id] = session

    if not client:
        raise HTTPException(status_code=500, detail="MISTRAL_API_KEY not configured")

    chunks  = session.get("chunks", [])
    use_rag = any(c.get("embedding") for c in chunks)
    sources = []

    if use_rag:
        try:
            q_resp = client.embeddings.create(
                model="mistral-embed",
                inputs=[body.message],
            )
            question_embedding = q_resp.data[0].embedding
            relevant       = find_relevant_chunks(question_embedding, chunks, top_k=5)
            evidence       = build_evidence_pack(relevant)
            selected_mode  = body.mode if body.mode in MODE_PROMPTS else "explication"
            mode_prompt    = MODE_PROMPTS[selected_mode]
            system_content = mode_prompt.format(evidence_pack=evidence)
            print(f"[lib_chat] Mode: {selected_mode} | RAG chunks: {len(relevant)}")
            sources = [
                {
                    "doc":   c["doc_name"],
                    "chunk": c["chunk_index"],
                    "text":  c["content"][:300],
                    "words": len(c["content"].split()),
                    "score": round(c.pop("_score", 0.0), 3),
                }
                for c in relevant
            ]
        except Exception as e:
            print(f"[lib_chat] RAG failed ({e}) — using master_md fallback")
            use_rag = False

    if not use_rag:
        selected_mode  = body.mode if body.mode in MODE_PROMPTS else "explication"
        mode_prompt    = MODE_PROMPTS[selected_mode]
        system_content = mode_prompt.format(
            evidence_pack=f"SOURCES:\n{session['master_md'][:6000]}"
        )

    messages_to_send = [{"role": "system", "content": system_content}]
    for h in body.history:
        if h.get("role") in ("user", "assistant") and h.get("content"):
            messages_to_send.append({"role": h["role"], "content": h["content"]})
    messages_to_send.append({"role": "user", "content": body.message})

    def generate():
        full_answer = ""
        try:
            if sources:
                yield f"data: {json.dumps({'sources': sources})}\n\n"

            stream = client.chat.stream(model="mistral-small-latest", messages=messages_to_send)
            for event in stream:
                delta = None
                try:
                    choice = event.data.choices[0]
                    if choice.delta and choice.delta.content:
                        delta = choice.delta.content
                except (AttributeError, IndexError, TypeError):
                    pass
                if delta:
                    full_answer += delta
                    yield f"data: {json.dumps({'token': delta})}\n\n"

            lib_save_student_message(session_id, "user", body.message)
            lib_save_student_message(session_id, "assistant", full_answer)
            if "messages" not in session:
                session["messages"] = []
            session["messages"].append({"role": "user",      "content": body.message})
            session["messages"].append({"role": "assistant",  "content": full_answer})

            yield "data: [DONE]\n\n"
        except Exception as e:
            yield f"data: {json.dumps({'error': str(e)})}\n\n"
            yield "data: [DONE]\n\n"

    return StreamingResponse(
        generate(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


@app.get("/api/student/dashboard/{session_id}")
def student_dashboard(session_id: str):
    """
    Returns all data needed for the student dashboard.
    Combines: morning briefing text, priority concepts,
    weekly progress, mastery by course.
    """
    db = get_db()

    session = get_session_with_fallback(session_id)
    if not session:
        raise HTTPException(status_code=404, detail="Session not found")

    messages = []
    try:
        if db:
            res = db.table("student_messages")\
                    .select("role,content,created_at")\
                    .eq("session_id", session_id)\
                    .order("created_at")\
                    .execute()
            messages = res.data or []
    except Exception:
        messages = session.get("messages", [])

    today = datetime.now(timezone.utc)
    week_days = []
    for i in range(6, -1, -1):
        day = today - timedelta(days=i)
        day_label = ["Lun", "Mar", "Mer", "Jeu", "Ven", "Sam", "Dim"][day.weekday()]
        day_str = day.strftime("%Y-%m-%d")
        count = sum(
            1 for m in messages
            if m.get("created_at", "")[:10] == day_str
            and m.get("role") == "user"
        )
        week_days.append({"day": day_label, "count": count, "date": day_str})

    week_total = sum(d["count"] for d in week_days)
    sessions_done = sum(1 for d in week_days if d["count"] > 0)

    concepts = []
    try:
        if db:
            res = db.table("concept_mastery")\
                    .select("*")\
                    .eq("session_id", session_id)\
                    .execute()
            concepts = res.data or []
    except Exception:
        concepts = []

    priority_concepts = []
    for c in concepts:
        try:
            last_reviewed = datetime.fromisoformat(
                c.get("last_reviewed", str(today)).replace("Z", "+00:00")
            )
            days_since = max(1, (today - last_reviewed).days)
            avg_mastery = max(0.01, (
                c.get("memorisation", 0) +
                c.get("comprehension", 0) +
                c.get("application", 0)
            ) / 3)
            score = (c.get("exam_weight", 1) * max(1, c.get("error_frequency", 1))) \
                    / (days_since * avg_mastery)
            priority_concepts.append({**c, "priority_score": round(score, 2)})
        except Exception:
            continue

    priority_concepts.sort(key=lambda x: x["priority_score"], reverse=True)
    top_concepts = priority_concepts[:3]

    courses = {}
    for c in concepts:
        course = c.get("course", "Cours général")
        if course not in courses:
            courses[course] = {"name": course, "concepts": [], "avg_mastery": 0}
        courses[course]["concepts"].append(c)

    course_list = []
    for course_name, data in courses.items():
        if data["concepts"]:
            avg = sum(
                (c.get("memorisation", 0) + c.get("comprehension", 0) + c.get("application", 0)) / 3
                for c in data["concepts"]
            ) / len(data["concepts"])
            course_list.append({
                "name":           course_name,
                "concept_count":  len(data["concepts"]),
                "avg_mastery":    round(avg, 2),
                "memorisation":   round(sum(c.get("memorisation", 0) for c in data["concepts"]) / len(data["concepts"]), 2),
                "comprehension":  round(sum(c.get("comprehension", 0) for c in data["concepts"]) / len(data["concepts"]), 2),
                "application":    round(sum(c.get("application", 0) for c in data["concepts"]) / len(data["concepts"]), 2),
            })

    doc_names = session.get("files", [])
    recent_concepts = [c.get("concept", "") for c in concepts[-3:]] if concepts else []

    return {
        "week_days":        week_days,
        "sessions_done":    sessions_done,
        "week_total":       week_total,
        "top_concepts":     top_concepts,
        "courses":          course_list,
        "doc_count":        len(doc_names),
        "message_count":    len([m for m in messages if m.get("role") == "user"]),
        "recent_concepts":  recent_concepts,
        "has_data":         len(concepts) > 0,
    }


@app.post("/api/student/morning-briefing/{session_id}")
async def morning_briefing(session_id: str):
    """
    Generate a personalised morning briefing message using Mistral.
    Returns a short motivating text with today's priority action.
    """
    if not client:
        return {"briefing": "Bonjour ! Prêt à travailler sur vos cours aujourd'hui ?"}

    try:
        dashboard = student_dashboard(session_id)
    except Exception:
        return {"briefing": "Bonjour ! Commencez votre session d'aujourd'hui."}

    top = dashboard.get("top_concepts", [])
    sessions_done = dashboard.get("sessions_done", 0)
    week_total = dashboard.get("week_total", 0)

    top_concept_text = ""
    if top:
        top_concept_text = f"Le concept prioritaire aujourd'hui est '{top[0].get('concept', '')}' ({top[0].get('course', '')})."

    prompt = f"""Tu es Kapsul, un tuteur pédagogique bienveillant et motivant.
Génère un message de briefing matinal COURT (2-3 phrases maximum) pour cet étudiant.

Données de l'étudiant cette semaine :
- Sessions de travail complétées : {sessions_done}/7 jours
- Messages échangés : {week_total}
- {top_concept_text}

RÈGLES :
- Commence toujours par "Bonjour !" ou une variante chaleureuse
- Mentionne un accomplissement concret si possible (ex: "Tu as consolidé X concepts hier")
- Termine par l'objectif du jour en une phrase claire et motivante
- Ton : bienveillant, encourageant, direct. Jamais anxiogène.
- Maximum 3 phrases. Pas de listes. Pas de titres.
- En français uniquement."""

    try:
        response = client.chat.complete(
            model="mistral-small-latest",
            messages=[{"role": "user", "content": prompt}],
        )
        briefing = response.choices[0].message.content or "Bonjour ! Prêt pour une nouvelle session ?"
        return {"briefing": briefing.strip()}
    except Exception as e:
        print(f"[briefing] Failed: {e}")
        return {"briefing": "Bonjour ! Votre session d'aujourd'hui vous attend. Bonne révision ! 🎯"}
